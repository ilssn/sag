"""上传文件 → Markdown 的路由、缓存与 302 MinerU 适配。"""

from __future__ import annotations

import asyncio
import io
import zipfile
from pathlib import Path
from types import SimpleNamespace
from typing import Any

import httpx
import pytest

from sag_api.core.config import Settings
from sag_api.core.errors import (
    ConfigurationError,
    ServiceUnavailableError,
    UpstreamError,
)
from sag_api.parsing import service
from sag_api.parsing.mineru import MinerUClient, _assert_public_host, _interpret_poll_payload
from sag_api.parsing.service import PreparedDocument
from sag_api.sag.dto import ProcessCheckpoint, ProcessOutcome


def _settings(**overrides: Any) -> Settings:
    return Settings(
        _env_file=None,
        data_dir="/tmp/sag-test-engine",
        upload_dir="/tmp/sag-test-uploads",
        **overrides,
    )


@pytest.mark.asyncio
async def test_parser_routes_markdown_and_markitdown_with_cache(tmp_path, monkeypatch):
    markdown = tmp_path / "already.md"
    markdown.write_text("# Already\n", encoding="utf-8")
    direct = await service.prepare_document(str(markdown), _settings())
    assert direct.path == str(markdown)
    assert direct.provider == "original"

    source = tmp_path / "notes.txt"
    source.write_text("hello", encoding="utf-8")
    calls: list[str] = []

    def convert(path: str) -> str:
        calls.append(path)
        return "# Converted\n\nhello"

    monkeypatch.setattr(service, "_markitdown_sync", convert)
    first = await service.prepare_document(str(source), _settings())
    second = await service.prepare_document(str(source), _settings())

    assert first.provider == "markitdown" and first.path.endswith(".parsed.markitdown.md")
    assert Path(first.path).read_text(encoding="utf-8") == "# Converted\n\nhello\n"
    assert second.cached is True and second.path == first.path
    assert calls == [str(source)]


@pytest.mark.asyncio
async def test_only_pdf_uses_configured_mineru(tmp_path, monkeypatch):
    settings = _settings(
        document_parser="auto",
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    seen: list[str] = []

    class FakeMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            seen.append(path)
            return "# From MinerU\n"

    monkeypatch.setattr(service, "MinerUClient", FakeMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# From MarkItDown\n")

    pdf = tmp_path / "paper.pdf"
    pdf.write_bytes(b"%PDF-fake")
    docx = tmp_path / "paper.docx"
    docx.write_bytes(b"fake-office")

    parsed_pdf = await service.prepare_document(str(pdf), settings)
    parsed_docx = await service.prepare_document(str(docx), settings)
    assert parsed_pdf.provider == "mineru"
    assert parsed_docx.provider == "markitdown"
    assert seen == [str(pdf)]


@pytest.mark.asyncio
async def test_pdf_without_complete_mineru_config_falls_back_to_markitdown(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# Local PDF\n")

    parsed = await service.prepare_document(
        str(source),
        _settings(document_parser="mineru", mineru_base_url="https://api.302.ai"),
    )
    assert parsed.provider == "markitdown"


@pytest.mark.asyncio
async def test_mineru_failure_falls_back_to_markitdown_and_reuses_cache(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0
    markitdown_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            nonlocal mineru_calls
            mineru_calls += 1
            if on_state:
                await on_state({**(state or {}), "task_id": "task-unavailable"})
            raise ServiceUnavailableError("No available models currently")

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# Local fallback\n"

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    states: list[dict[str, Any]] = []

    first = await service.prepare_document(
        str(source), settings, on_state=lambda state: _record(states, state)
    )
    second = await service.prepare_document(
        str(source),
        settings,
        state=states[-1],
        on_state=lambda state: _record(states, state),
    )

    assert first.provider == "markitdown" and first.path.endswith(".parsed.markitdown.md")
    assert Path(first.path).read_text(encoding="utf-8") == "# Local fallback\n"
    assert second.path == first.path and second.cached is True
    assert mineru_calls == 1 and markitdown_calls == 1
    assert states[-1]["provider"] == "mineru"
    assert states[-1]["status"] == "fallback_done"
    assert states[-1]["fallback"]["provider"] == "markitdown"
    assert states[-1]["fallback"]["mineru_error"] == "No available models currently"
    assert states[-1]["fallback"]["status"] == "done"
    assert first.fallback_from == "mineru"
    assert first.fallback_error == "No available models currently"


@pytest.mark.asyncio
async def test_document_fails_only_when_mineru_and_markitdown_both_fail(
    tmp_path, monkeypatch
):
    source = tmp_path / "broken.pdf"
    source.write_bytes(b"%PDF-broken")

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            raise ServiceUnavailableError("remote parser failed")

    def fail_locally(_path: str) -> str:
        raise RuntimeError("local parser failed")

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", fail_locally)
    states: list[dict[str, Any]] = []
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    with pytest.raises(ServiceUnavailableError, match="MinerU.*MarkItDown"):
        await service.prepare_document(
            str(source), settings, on_state=lambda state: _record(states, state)
        )
    assert states[-1]["status"] == "fallback_failed"
    assert states[-1]["fallback"]["markitdown_error"].endswith("local parser failed")


@pytest.mark.asyncio
async def test_mineru_state_callback_failure_does_not_trigger_markitdown(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    markitdown_calls = 0

    class MinerUWithState:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            assert on_state is not None
            await on_state({**(state or {}), "task_id": "task-1"})
            return "# Never reached\n"

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# Should not run\n"

    state_writes = 0

    async def fail_to_persist(_state: dict[str, Any]) -> None:
        nonlocal state_writes
        state_writes += 1
        if state_writes == 2:
            raise RuntimeError("database commit failed")

    monkeypatch.setattr(service, "MinerUClient", MinerUWithState)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    with pytest.raises(RuntimeError, match="database commit failed"):
        await service.prepare_document(
            str(source), settings, on_state=fail_to_persist
        )
    assert state_writes == 2
    assert markitdown_calls == 0


@pytest.mark.asyncio
async def test_changed_mineru_config_retries_remote_after_cached_fallback(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            nonlocal mineru_calls
            mineru_calls += 1
            raise UpstreamError("remote unavailable")

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", lambda path: "# Local fallback\n")
    states: list[dict[str, Any]] = []
    first_settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-first",
    )
    second_settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-changed",
    )

    await service.prepare_document(
        str(source), first_settings, on_state=lambda state: _record(states, state)
    )
    previous_state = states[-1]
    await service.prepare_document(
        str(source),
        second_settings,
        state=previous_state,
        on_state=lambda state: _record(states, state),
    )

    assert mineru_calls == 2
    assert states[-1]["key_fingerprint"] != previous_state["key_fingerprint"]


@pytest.mark.asyncio
async def test_concurrent_pdf_parsing_creates_only_one_mineru_task(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    started = asyncio.Event()
    release = asyncio.Event()
    calls = 0

    class FakeMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            nonlocal calls
            calls += 1
            started.set()
            await release.wait()
            return "# Parsed once\n"

    monkeypatch.setattr(service, "MinerUClient", FakeMinerU)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )
    first = asyncio.create_task(service.prepare_document(str(source), settings))
    second = asyncio.create_task(service.prepare_document(str(source), settings))
    await started.wait()
    release.set()
    results = await asyncio.gather(first, second)

    assert calls == 1
    assert results[0].path == results[1].path
    assert sorted(result.cached for result in results) == [False, True]


@pytest.mark.asyncio
async def test_concurrent_mineru_failure_creates_one_task_and_one_fallback(
    tmp_path, monkeypatch
):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    mineru_calls = 0
    markitdown_calls = 0

    class FailingMinerU:
        def __init__(self, _settings):
            pass

        async def parse(self, path, *, state=None, on_state=None):
            nonlocal mineru_calls
            mineru_calls += 1
            await asyncio.sleep(0.01)
            raise ServiceUnavailableError("No available models currently")

    def convert(_path: str) -> str:
        nonlocal markitdown_calls
        markitdown_calls += 1
        return "# One local fallback\n"

    monkeypatch.setattr(service, "MinerUClient", FailingMinerU)
    monkeypatch.setattr(service, "_markitdown_sync", convert)
    settings = _settings(
        mineru_base_url="https://api.302.ai",
        mineru_api_key="sk-test",
    )

    results = await asyncio.gather(
        service.prepare_document(str(source), settings),
        service.prepare_document(str(source), settings),
    )

    assert mineru_calls == 1 and markitdown_calls == 1
    assert results[0].path == results[1].path
    assert all(result.provider == "markitdown" for result in results)
    assert sorted(result.cached for result in results) == [False, True]


@pytest.mark.asyncio
async def test_document_job_sends_parsed_markdown_to_engine(monkeypatch):
    from sag_api.db.models import Document, Source
    from sag_api.jobs import tasks

    document = SimpleNamespace(
        id="doc-1",
        source_id="source-1",
        storage_path="/uploads/original.pdf",
        status=None,
        error=None,
        chunk_count=0,
        event_count=0,
        progress=0,
        token_usage=0,
        sag_source_id=None,
    )
    source = SimpleNamespace(
        id="source-1",
        sag_source_config_id="sag-source-1",
        chunk_count=0,
        event_count=0,
    )
    job = SimpleNamespace(id="job-1", document_id="doc-1", progress=0.0, payload={})

    class FakeSession:
        async def get(self, model, _id):
            return document if model is Document else source if model is Source else None

        async def commit(self):
            pass

        async def execute(self, _statement):
            pass

        async def refresh(self, _instance, attribute_names=None):
            pass

    prepared_calls: list[str] = []

    async def fake_prepare(path, settings, *, state=None, on_state=None):
        prepared_calls.append(path)
        return PreparedDocument("/uploads/original.pdf.parsed.markitdown.md", "markitdown")

    class FakeEngineManager:
        seen_path = ""

        async def process_document(
            self,
            source_config_id,
            path,
            *,
            source,
            on_stage,
            checkpoint,
            on_checkpoint,
            should_pause,
            max_concurrency,
        ):
            self.seen_path = path
            assert max_concurrency == tasks.settings.document_extract_concurrency
            await on_stage("loading")
            await on_checkpoint(
                ProcessCheckpoint(
                    source_id="engine-doc",
                    chunk_ids=["chunk-1", "chunk-2"],
                    processed_chunk_ids=["chunk-1"],
                    event_count=1,
                    event_ids=["event-1"],
                    token_usage=1234,
                )
            )
            await on_stage("extracting")
            return ProcessOutcome(
                source_id="engine-doc",
                chunk_count=2,
                event_count=1,
                chunk_ids=["chunk-1", "chunk-2"],
                processed_chunk_ids=["chunk-1", "chunk-2"],
                token_usage=2468,
            )

    monkeypatch.setattr(tasks, "prepare_document", fake_prepare)
    engine = FakeEngineManager()
    await tasks.process_document(FakeSession(), job, engine_manager=engine)

    assert prepared_calls == ["/uploads/original.pdf"]
    assert engine.seen_path.endswith(".md")
    assert document.status.value == "ready"
    assert document.chunk_count == 2 and document.event_count == 1
    assert document.progress == 100 and document.token_usage == 2468


class _FakeAsyncClient:
    responses: list[httpx.Response] = []
    calls: list[tuple[str, str, dict[str, Any]]] = []

    def __init__(self, *args, **kwargs):
        pass

    async def __aenter__(self):
        return self

    async def __aexit__(self, exc_type, exc, tb):
        return False

    @classmethod
    def reset(cls, responses: list[httpx.Response]) -> None:
        cls.responses = list(responses)
        cls.calls = []

    @classmethod
    def _next(cls) -> httpx.Response:
        assert cls.responses, "unexpected HTTP request"
        return cls.responses.pop(0)

    async def post(self, url: str, **kwargs):
        self.calls.append(("POST", url, kwargs))
        return self._next()

    async def request(self, method: str, url: str, **kwargs):
        self.calls.append((method, url, kwargs))
        return self._next()

    class _Stream:
        def __init__(self, response: httpx.Response):
            self.response = response

        async def __aenter__(self):
            return self.response

        async def __aexit__(self, exc_type, exc, tb):
            return False

    def stream(self, method: str, url: str, **kwargs):
        self.calls.append(("DOWNLOAD", url, kwargs))
        return self._Stream(self._next())


def _response(
    *,
    json: Any = None,
    content: bytes | None = None,
    content_type: str = "application/json",
    url: str = "https://api.302.ai/test",
    status: int = 200,
) -> httpx.Response:
    request = httpx.Request("GET", url)
    if content is not None:
        return httpx.Response(
            status, content=content, headers={"content-type": content_type}, request=request
        )
    return httpx.Response(status, json=json, request=request)


def _result_zip(markdown: str) -> bytes:
    target = io.BytesIO()
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("images/ignored.txt", "x")
        archive.writestr("full.md", markdown)
    return target.getvalue()


def _simple_pdf(text: str) -> bytes:
    """生成带可提取文本层的最小 PDF，避免测试依赖 PDF 写入库。"""
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, 1):
        offsets.append(len(output))
        output.extend(f"{number} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode())
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    return bytes(output)


def _simple_docx(path: Path, text: str) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/word/document.xml"
                ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1"
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"
                Target="word/document.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
              <w:body><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:body>
            </w:document>""",
        )


def test_real_markitdown_converts_pdf_and_office_files(tmp_path):
    """依赖安装烟测：核心格式确实能产出可供引擎摄取的 Markdown。"""
    from openpyxl import Workbook
    from pptx import Presentation

    from sag_api.parsing.service import _markitdown_sync

    pdf = tmp_path / "sample.pdf"
    pdf.write_bytes(_simple_pdf("Muse PDF marker"))

    docx = tmp_path / "sample.docx"
    _simple_docx(docx, "Muse DOCX marker")

    pptx = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Muse PPTX marker"
    presentation.save(pptx)

    xlsx = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "Muse XLSX marker"
    workbook.save(xlsx)

    assert "Muse PDF marker" in _markitdown_sync(str(pdf))
    assert "Muse DOCX marker" in _markitdown_sync(str(docx))
    assert "Muse PPTX marker" in _markitdown_sync(str(pptx))
    assert "Muse XLSX marker" in _markitdown_sync(str(xlsx))


@pytest.mark.asyncio
async def test_mineru_upload_create_poll_and_download_zip(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset(
        [
            _response(
                json={"code": 200, "data": "https://file.302.ai/input.pdf", "message": "success"}
            ),
            _response(json="task-123"),
            _response(json="processing"),
            _response(json='{"status":"done","full_zip_url":"https://file.302.ai/result.zip"}'),
            _response(
                content=_result_zip("# Parsed\n\nMinerU result"),
                content_type="application/zip",
                url="https://file.302.ai/result.zip",
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    states: list[dict[str, Any]] = []
    client = MinerUClient(
        _settings(
            mineru_base_url="https://api.302.ai",
            mineru_api_key="sk-mineru",
            mineru_poll_interval=0.001,
            mineru_poll_timeout=1,
        )
    )

    markdown = await client.parse(str(source), on_state=lambda state: _record(states, state))

    assert markdown == "# Parsed\n\nMinerU result\n"
    assert any(state.get("task_id") == "task-123" for state in states)
    assert [call[0] for call in _FakeAsyncClient.calls] == [
        "POST",
        "POST",
        "GET",
        "GET",
        "DOWNLOAD",
    ]
    assert _FakeAsyncClient.calls[0][1] == "https://api.302.ai/302/upload-file"
    assert _FakeAsyncClient.calls[1][2]["json"]["version"] == "2.5"
    assert "Authorization" not in _FakeAsyncClient.calls[-1][2]


async def _record(target: list[dict[str, Any]], state: dict[str, Any]) -> None:
    target.append(dict(state))


@pytest.mark.asyncio
async def test_mineru_retry_reuses_upload_and_task_id(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset([_response(json="# Ready\n\nNo duplicate paid task")])
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(
        _settings(
            mineru_base_url="https://api.302.ai",
            mineru_api_key="sk-mineru",
            mineru_poll_interval=0.001,
            mineru_poll_timeout=1,
        )
    )

    markdown = await client.parse(
        str(source),
        state={"upload_url": "https://file.302.ai/input.pdf", "task_id": "existing-task"},
    )
    assert markdown.startswith("# Ready")
    assert [call[0] for call in _FakeAsyncClient.calls] == ["GET"]


@pytest.mark.asyncio
async def test_mineru_accepts_immediate_result_url_from_create(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset(
        [
            _response(json={"code": 200, "data": "https://file.302.ai/input.pdf"}),
            _response(json="https://file.302.ai/result.zip"),
            _response(
                content=_result_zip("# Immediate result"),
                content_type="application/zip",
                url="https://file.302.ai/result.zip",
            ),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(
        _settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-mineru")
    )

    assert await client.parse(str(source)) == "# Immediate result\n"
    assert [call[0] for call in _FakeAsyncClient.calls] == ["POST", "POST", "DOWNLOAD"]


@pytest.mark.asyncio
async def test_mineru_result_download_maps_bounded_error_body(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    oversized_error = {"data": {"err_msg": "result temporarily unavailable"}}
    _FakeAsyncClient.reset(
        [
            _response(json={"code": 200, "data": "https://file.302.ai/input.pdf"}),
            _response(json="https://file.302.ai/result.md"),
            _response(json=oversized_error, status=503),
        ]
    )
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(
        _settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-mineru")
    )

    with pytest.raises(ServiceUnavailableError, match="result temporarily unavailable"):
        await client.parse(str(source))


@pytest.mark.asyncio
async def test_mineru_auth_error_is_configuration_error(tmp_path, monkeypatch):
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    _FakeAsyncClient.reset([_response(json={"message": "bad key"}, status=401)])
    monkeypatch.setattr(httpx, "AsyncClient", _FakeAsyncClient)
    client = MinerUClient(
        _settings(mineru_base_url="https://api.302.ai", mineru_api_key="sk-bad")
    )
    with pytest.raises(ConfigurationError, match="API Key"):
        await client.parse(str(source))


def test_mineru_zip_requires_markdown():
    target = io.BytesIO()
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("result.json", "{}")
    from sag_api.parsing.mineru import _markdown_from_zip

    with pytest.raises(UpstreamError, match="没有 Markdown"):
        _markdown_from_zip(target.getvalue(), 1024)


def test_mineru_pending_and_nested_failure_payloads_are_not_misclassified():
    pending = {
        "status": "processing",
        "data": {"url": "https://file.302.ai/input.pdf"},
    }
    assert _interpret_poll_payload(pending, "task-1")[0] == "pending"

    nested_failure = {
        "code": 200,
        "msg": "ok",
        "data": {"status": "failed", "err_msg": "bad pdf"},
    }
    assert _interpret_poll_payload(nested_failure, "task-1") == ("failed", "bad pdf")
    assert _interpret_poll_payload("A0202", "task-1")[0] == "failed"


@pytest.mark.asyncio
async def test_mineru_result_download_rejects_private_hosts():
    with pytest.raises(UpstreamError, match="内网"):
        await _assert_public_host("127.0.0.1", 80)
    with pytest.raises(UpstreamError, match="内网"):
        await _assert_public_host("169.254.169.254", 80)
